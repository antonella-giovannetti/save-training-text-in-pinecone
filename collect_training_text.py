import os
from pptx import Presentation
from transformers import AutoTokenizer,  AutoModel
from sentence_transformers import SentenceTransformer
import pinecone

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def chunk_text(text, chunk_size, chunk_overlap):
    tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/paraphrase-MiniLM-L6-v2")
    tokens = tokenizer.encode(text, add_special_tokens=False) 
    chunks = []
    
    for i in range(0, len(tokens), chunk_size - chunk_overlap):
        chunk = tokens[i:i + chunk_size]
        chunk_text = tokenizer.decode(chunk)
        chunks.append(chunk_text)
    
    return chunks

chunk_size = 384
chunk_overlap = 50

def process_pptx_files(folder_path, chunk_size, chunk_overlap):
    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx'):
            file_path = os.path.join(folder_path, filename)
            print(f"Processing file: {filename}")
            text = extract_text_from_pptx(file_path)
            chunks = chunk_text(text, chunk_size, chunk_overlap)
            print(f"Text extracted from {filename} and chunked:")
            for i, chunk in enumerate(chunks):
                print(f"Chunk {i+1}:")
                print(chunk)
                print("-" * 50)

folder_path = 'data/pptx'
chunk_size = 384
chunk_overlap = 50
process_pptx_files(folder_path, chunk_size, chunk_overlap)

def create_embeddings(chunks):
    model = SentenceTransformer("sentence-transformers/paraphrase-MiniLM-L6-v2")
    embeddings = model.encode(chunks)
    return embeddings

def store_in_pinecone(embeddings, chunks, filename):
    pinecone.init(api_key="your-api-key", environment="your-environment")
    index = pinecone.Index("your-index-name")
    
    for i, (embedding, chunk) in enumerate(zip(embeddings, chunks)):
        index.upsert([(f"{filename}_chunk_{i}", embedding.tolist(), {"text": chunk})])

def process_pptx_files(folder_path, chunk_size, chunk_overlap):
    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx'):
            file_path = os.path.join(folder_path, filename)
            print(f"Processing file: {filename}")
            
            # Extract text
            text = extract_text_from_pptx(file_path)
            
            # Chunk the text
            chunks = chunk_text(text, chunk_size, chunk_overlap)
            
            print(f"Text extracted from {filename} and chunked into {len(chunks)} chunks.")
            
            # Create embeddings
            embeddings = create_embeddings(chunks)
            print(f"Created {len(embeddings)} embeddings.")
            
            # Store in Pinecone
            store_in_pinecone(embeddings, chunks, filename)
            print(f"Stored embeddings and chunks in Pinecone for {filename}")
            
            print("-" * 50)

# Usage
folder_path = 'data/pptx'
chunk_size = 384
chunk_overlap = 50
process_pptx_files(folder_path, chunk_size, chunk_overlap)

